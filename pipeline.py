
import shutil, subprocess, json, re, copy, io, zipfile
from pathlib import Path

import openpyxl
from pptx import Presentation
from pptx.util import Pt
from pptx.oxml.ns import qn
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
import numpy as np

SCRIPTS_DIR = Path(__file__).parent / "scripts"

NIVELES = {
    "marca":        "Análisis por Marca",
    "segmento":     "Análisis por Segmento",
    "subcategoria": "Análisis por Subcategoría",
}

def short(name):    return (name or "").replace("Total ", "").strip()
def fmt_m(v):
    if v is None: return "–"
    av = abs(v)
    if av >= 1e9:  return f"${v/1e9:.1f}B"
    if av >= 1e6:  return f"${v/1e6:.1f}M"
    if av >= 1e3:  return f"${v/1e3:.0f}K"
    return f"${v:,.0f}"
def fmt_pct(v):
    return "–" if v is None else f"{v*100:+.1f}%"

def inject_data(input_path, work_xlsx):
    wb_in  = openpyxl.load_workbook(input_path)
    ws_in  = wb_in["Principal"]
    wb_tpl = openpyxl.load_workbook(work_xlsx)
    ws_tpl = wb_tpl["02_Marca_OUTPUT"]
    for r in range(1, ws_in.max_row + 1):
        for c in range(1, ws_in.max_column + 1):
            ws_tpl.cell(r, c).value = ws_in.cell(r, c).value
    wb_tpl.save(work_xlsx)

def recalc(work_xlsx):
    subprocess.run(
        ["python3", str(SCRIPTS_DIR / "recalc.py"), str(work_xlsx), "90"],
        capture_output=True,
        cwd=str(SCRIPTS_DIR)
    )

def read_data(work_xlsx):
    wb = openpyxl.load_workbook(work_xlsx, data_only=True)
    ws = wb["02_Marca_OUTPUT"]
    g  = lambda r, c: ws.cell(r, c).value

    t1 = []
    for r in range(21, 32):
        brand = g(r, 2)
        if not brand or brand == "Total": continue
        t1.append({"brand": brand, "ventas": g(r,3), "unidades": g(r,4),
                   "clientes": g(r,5), "var_v": g(r,9), "var_u": g(r,10), "var_c": g(r,11)})

    sv, su = [], []
    for r in range(21, 32):
        brand = g(r, 16)
        if not brand or brand == "Total": continue
        sv.append({"brand": brand, "sh_ant": g(r,19), "sh_rec": g(r,20), "dif_sh": g(r,22)})
        su.append({"brand": g(r,24), "sh_ant_u": g(r,27), "sh_rec_u": g(r,28), "dif_sh_u": g(r,30)})

    mot = []
    for r in range(39, 50):
        brand = g(r, 2)
        if not brand: continue
        mot.append({"brand": brand, "penetracion": g(r,3), "frecuencia": g(r,4),
                    "precio": g(r,5), "clientes_s": g(r,6), "unid_visita": g(r,7)})

    return {"t1": t1, "sv": sv, "su": su, "mot": mot}

W = dict(facecolor="white")

def fig_motivadores(data, out):
    mot = data["mot"]
    if not mot: return None
    labels = [short(m["brand"]) for m in mot]
    keys   = ["penetracion","frecuencia","precio","clientes_s","unid_visita"]
    cats   = ["Penetración","Frecuencia","Precio Prom.","Clientes Super.","Unid×Visita"]
    colors = ["#4472C4","#ED7D31","#A9D18E","#FFC000","#5B9BD5"]
    matrix = np.array([[m[k] or 0 for k in keys] for m in mot], float) / 1e6

    fig, ax = plt.subplots(figsize=(9.0, 6.3), **W)
    fig.patch.set_facecolor("white"); ax.set_facecolor("white")
    x = np.arange(len(labels)); bp, bn = np.zeros(len(labels)), np.zeros(len(labels))
    for i, (cat, col) in enumerate(zip(cats, colors)):
        vals = matrix[:, i]
        p, n = np.where(vals>=0, vals, 0), np.where(vals<0, vals, 0)
        ax.bar(x, p, bottom=bp, color=col, label=cat, width=0.65, edgecolor="white", lw=0.4)
        ax.bar(x, n, bottom=bn, color=col, width=0.65, edgecolor="white", lw=0.4)
        bp += p; bn += n
    ax.axhline(0, color="#555", lw=0.8)
    ax.set_xticks(x); ax.set_xticklabels(labels, rotation=40, ha="right", fontsize=9)
    ax.set_ylabel("Millones COP", fontsize=10)
    ax.set_title("Motivadores del Cambio en Ventas", fontsize=12, fontweight="bold", pad=10)
    ax.legend(loc="upper right", fontsize=8, framealpha=0.95)
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda v,_: f"${v:.0f}M"))
    ax.spines[["top","right"]].set_visible(False)
    plt.tight_layout(pad=1.2)
    fig.savefig(out, dpi=150, bbox_inches="tight", facecolor="white"); plt.close(fig)
    return out

def fig_share(data, out, mode="ventas"):
    items = data["sv"] if mode=="ventas" else [x for x in data["su"] if x.get("brand")]
    if not items: return None
    labels = [short(x["brand"]) for x in items]
    if mode == "ventas":
        sa = [(x["sh_ant"] or 0)*100 for x in items]
        sr = [(x["sh_rec"] or 0)*100 for x in items]
        df = [x["dif_sh"] or 0 for x in items]
        title, ca, cr = "Share en Ventas $ (EPOS)", "#BDD7EE", "#2E75B6"
    else:
        sa = [(x["sh_ant_u"] or 0)*100 for x in items]
        sr = [(x["sh_rec_u"] or 0)*100 for x in items]
        df = [x["dif_sh_u"] or 0 for x in items]
        title, ca, cr = "Share en Unidades (EPOS)", "#C6EFCE", "#548235"

    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(11.0, 5.8),
                                   gridspec_kw={"width_ratios":[2.8,1]}, **W)
    fig.patch.set_facecolor("white")
    x = np.arange(len(labels)); w = 0.38
    ax1.bar(x-w/2, sa, w, label="Anterior", color=ca, edgecolor="white")
    ax1.bar(x+w/2, sr, w, label="Reciente", color=cr, edgecolor="white")
    ax1.set_xticks(x); ax1.set_xticklabels(labels, rotation=40, ha="right", fontsize=9)
    ax1.set_ylabel("Share (%)", fontsize=10); ax1.set_title(title, fontsize=12, fontweight="bold")
    ax1.legend(fontsize=8.5); ax1.spines[["top","right"]].set_visible(False)
    ax1.yaxis.set_major_formatter(mticker.FuncFormatter(lambda v,_: f"{v:.1f}%"))
    cols2 = ["#2ECC71" if d>=0 else "#E74C3C" for d in df]
    ax2.barh(labels, df, color=cols2, edgecolor="white", height=0.6)
    ax2.axvline(0, color="#555", lw=0.8)
    ax2.set_title("Δ Share (pp)", fontsize=11, fontweight="bold")
    ax2.spines[["top","right"]].set_visible(False); ax2.tick_params(axis="y", labelsize=9)
    offset = max(abs(d) for d in df)*0.04 if any(df) else 0.05
    for i, d in enumerate(df):
        ax2.text(d+(offset if d>=0 else -offset), i, f"{d:+.1f}",
                 va="center", ha="left" if d>=0 else "right", fontsize=8)
    plt.tight_layout(pad=1.2)
    fig.savefig(out, dpi=150, bbox_inches="tight", facecolor="white"); plt.close(fig)
    return out

def fig_tablas_combinadas(data, out):
    rows = data["t1"]; mot = data["mot"]
    if not rows and not mot: return None
    cols1 = ["Nombre","Ventas","Unidades","Clientes","Δ% Ventas","Δ% Unid.","Δ% Cli."]
    td1 = [[short(r["brand"]), fmt_m(r["ventas"]),
            f"{int(r['unidades']):,}" if r["unidades"] else "–",
            f"{int(r['clientes']):,}" if r["clientes"] else "–",
            fmt_pct(r["var_v"]), fmt_pct(r["var_u"]), fmt_pct(r["var_c"])] for r in rows]
    cols2 = ["Nombre","Penetración","Frecuencia","Precio","Cli.Super.","Unid×Vis."]
    td2 = [[short(m["brand"]), fmt_m(m["penetracion"]), fmt_m(m["frecuencia"]),
            fmt_m(m["precio"]), fmt_m(m["clientes_s"]), fmt_m(m["unid_visita"])] for m in mot]

    fig, (ax1, ax2) = plt.subplots(2, 1, figsize=(8.5, 9.5), facecolor="white",
                                   gridspec_kw={"height_ratios":[1,1], "hspace":0.6})
    fig.patch.set_facecolor("white")
    for ax, cols, td, title in [
        (ax1, cols1, td1, "Tabla 1 – Desempeño Comercial"),
        (ax2, cols2, td2, "Tabla 2 – Contribuciones al Cambio"),
    ]:
        ax.axis("off")
        tbl = ax.table(cellText=td, colLabels=cols, loc="center", cellLoc="center")
        tbl.auto_set_font_size(False); tbl.set_fontsize(7.5)
        tbl.auto_set_column_width(range(len(cols))); tbl.scale(1, 1.3)
        for c in range(len(cols)):
            tbl[0,c].set_facecolor("#2E75B6"); tbl[0,c].set_text_props(color="white", fontweight="bold")
        for ri, row in enumerate(td, 1):
            base = "#F2F7FF" if ri%2==0 else "white"
            for ci in range(len(cols)):
                cell = tbl[ri,ci]; val = row[ci]
                if ci > 0 and val != "–":
                    cell.set_facecolor("#FFE0E0" if val.startswith("-") else
                                       "#E2EFDA" if val.startswith("+") else base)
                else: cell.set_facecolor(base)
                if ci == 0: cell.get_text().set_ha("left")
        ax.set_title(title, fontsize=9.5, fontweight="bold", loc="left", x=0.0, pad=5)
    fig.savefig(out, dpi=150, bbox_inches="tight", facecolor="white"); plt.close(fig)
    return out

def ai_insights(data, periodo="", nivel="marca"):
    label = NIVELES.get(nivel, nivel)
    try:
        import urllib.request
        t1 = data["t1"]; sv = data["sv"]
        summary = f"Nivel: {label}\nPeríodo: {periodo}\nDesempeño:\n"
        for r in t1[:8]:
            summary += f"  {short(r['brand'])}: Ventas {fmt_m(r['ventas'])} ({fmt_pct(r['var_v'])}), Clientes {fmt_pct(r['var_c'])}\n"
        summary += "Share:\n"
        for s in sv[:8]:
            if s["dif_sh"]: summary += f"  {short(s['brand'])}: Δ Share {s['dif_sh']:+.1f}pp\n"

        payload = json.dumps({
            "model": "claude-sonnet-4-20250514", "max_tokens": 400,
            "messages": [{"role": "user", "content":
                f"""Analista consumo masivo Colombia. Genera DOS títulos de slides PowerPoint
para Lavalozas a nivel de {label}. Máximo 2 renglones, con los números más relevantes.
{summary}
Responde SOLO en JSON: {{"slide_motivadores": "...", "slide_share": "..."}}"""}]
        }).encode()
        req = urllib.request.Request("https://api.anthropic.com/v1/messages", data=payload,
            headers={"Content-Type":"application/json","anthropic-version":"2023-06-01"})
        with urllib.request.urlopen(req, timeout=25) as resp:
            text = re.sub(r"```json|```","", json.loads(resp.read())["content"][0]["text"]).strip()
            return json.loads(text)
    except:
        return {"slide_motivadores": f"[{label}] Motivadores clave del cambio en ventas",
                "slide_share":       f"[{label}] Evolución del share de ventas y unidades"}

def is_content_slide(slide):
    has_chart = any(s.shape_type == 3  for s in slide.shapes)
    has_table = any(s.shape_type == 19 for s in slide.shapes)
    no_edit   = any("No Editar" in (s.text_frame.text if s.has_text_frame else "")
                    for s in slide.shapes)
    return (has_chart or has_table) and not no_edit

def is_separator_slide(slide):
    return "Section" in slide.slide_layout.name and not is_content_slide(slide)

def replace_shape(slide, name, img_path):
    for shape in slide.shapes:
        if shape.name == name:
            l, t, w, h = shape.left, shape.top, shape.width, shape.height
            shape._element.getparent().remove(shape._element)
            slide.shapes.add_picture(str(img_path), l, t, w, h)
            return True
    return False

def _set_text_safe(slide, text, font_size=Pt(32)):
    for shape in slide.shapes:
        if shape.has_text_frame:
            tf = shape.text_frame
            try:
                run0 = tf.paragraphs[0].runs[0]
                size  = run0.font.size  or font_size
                bold  = run0.font.bold
                color = run0.font.color.rgb if run0.font.color.type else None
            except:
                size = font_size; bold = True; color = None
            tf.clear()
            p = tf.paragraphs[0]; p.text = text
            run = p.runs[0] if p.runs else p.add_run()
            run.font.size = size; run.font.bold = bold if bold is not None else True
            if color: run.font.color.rgb = color
            return

def _clone_slide(target_prs, source_slide):
    layout = source_slide.slide_layout
    new_slide = target_prs.slides.add_slide(layout)
    src_sp_tree = source_slide.shapes._spTree
    tgt_sp_tree = new_slide.shapes._spTree
    for child in list(tgt_sp_tree):
        tgt_sp_tree.remove(child)
    for child in src_sp_tree:
        tgt_sp_tree.append(copy.deepcopy(child))
    return new_slide

def build_multilevel_pptx(template_pptx, out_path, levels_data):
    src = Presentation(template_pptx)
    sep_src  = next((s for s in src.slides if is_separator_slide(s)), None)
    c_slides = [s for s in src.slides if is_content_slide(s)]
    c_tpl1   = c_slides[0] if len(c_slides) > 0 else None
    c_tpl2   = c_slides[1] if len(c_slides) > 1 else None
    last_src = src.slides[-1]

    out = Presentation(template_pptx)
    sldIdLst = out.slides._sldIdLst
    for sldId in list(sldIdLst)[1:]:
        sldIdLst.remove(sldId)

    for level_info in levels_data:
        nivel    = level_info["nivel"]
        imgs     = level_info["imgs"]
        insights = level_info["insights"]
        label    = NIVELES.get(nivel, nivel.capitalize())

        if sep_src:
            sep = _clone_slide(out, sep_src)
            for shape in sep.shapes:
                if shape.has_text_frame and "Text Placeholder" in shape.name:
                    _set_text_safe(sep, label)
                    break

        if c_tpl1:
            c1 = _clone_slide(out, c_tpl1)
            for name, path in imgs.items():
                replace_shape(c1, name, path)
            _set_text_safe(c1, insights[0] if insights else "", Pt(18))

        if c_tpl2:
            c2 = _clone_slide(out, c_tpl2)
            for name, path in imgs.items():
                replace_shape(c2, name, path)
            _set_text_safe(c2, insights[1] if len(insights)>1 else "", Pt(18))

    _clone_slide(out, last_src)

    # Deduplicate zip entries (python-pptx creates duplicates when cloning slides)
    buf = io.BytesIO()
    out.save(buf)
    buf.seek(0)
    clean = io.BytesIO()
    seen: dict = {}
    with zipfile.ZipFile(buf, "r") as src_zip, \
         zipfile.ZipFile(clean, "w", zipfile.ZIP_DEFLATED) as dst_zip:
        for info in src_zip.infolist():
            if info.filename not in seen:
                seen[info.filename] = True
                dst_zip.writestr(info, src_zip.read(info.filename))
    out_path.write_bytes(clean.getvalue())

def process_nivel(input_path, template_xlsx, workdir, nivel):
    work_xlsx = workdir / f"work_{nivel}.xlsx"
    shutil.copy(template_xlsx, work_xlsx)

    periodo = ""
    try:
        wb = openpyxl.load_workbook(input_path)
        ws = wb.get("Detalles del Pedido")
        if ws:
            for row in ws.iter_rows(values_only=True):
                if row and "Período" in str(row[1] or ""):
                    periodo = str(row[2] or ""); break
    except: pass

    inject_data(input_path, work_xlsx)
    recalc(work_xlsx)
    data = read_data(work_xlsx)

    imgs = {
        "Chart 3": fig_motivadores(data,       workdir/f"{nivel}_mot.png"),
        "Table 4": fig_tablas_combinadas(data,  workdir/f"{nivel}_tablas.png"),
        "Chart 4": fig_share(data,              workdir/f"{nivel}_sv.png", "ventas"),
        "Chart 1": fig_share(data,              workdir/f"{nivel}_su.png", "unidades"),
    }
    imgs = {k: v for k, v in imgs.items() if v}
    ins  = ai_insights(data, periodo, nivel)
    return imgs, [ins["slide_motivadores"], ins["slide_share"]], periodo

def run_pipeline(input_files, template_xlsx, template_pptx, workdir):
    workdir = Path(workdir)
    levels_data = []
    periodos = []

    for nivel in ["marca", "segmento", "subcategoria"]:
        if nivel not in input_files:
            continue
        imgs, insights, periodo = process_nivel(
            input_files[nivel], template_xlsx, workdir, nivel)
        levels_data.append({"nivel": nivel, "imgs": imgs, "insights": insights})
        if periodo: periodos.append(periodo)

    out_path = workdir / "output.pptx"
    build_multilevel_pptx(template_pptx, out_path, levels_data)

    return out_path.read_bytes(), f"{len(levels_data)} nivel(es) · {periodos[0] if periodos else ''}"
